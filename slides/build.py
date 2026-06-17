#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Build the group-meeting deck for paper_zh.tex (MZM/DPMZM affine bias-control).

Run order (see run.ps1 / commands):
  1. python gen_formulas.py        -> formulas.json (display + table + marked inline)
  2. python render_math_assets.py formulas.json --output-dir math --font-size 14
  3. python build.py               -> 组会汇报-仿射偏压控制.pptx
"""
import json, math, os, sys, tempfile
from pathlib import Path
from PIL import Image, ImageChops
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

_SKILL_DIR = os.environ.get("CLAUDE_SKILL_DIR") or os.environ.get("GROUP_MEETING_PPTX_SKILL_DIR")
if not _SKILL_DIR:
    _SKILL_DIR = r"C:\Users\ckdfs\.claude\skills\group-meeting-pptx"
sys.path.insert(0, str(Path(_SKILL_DIR) / "scripts"))
from deckbuilder import (TemplateDeck, fill_textframe, text_bottom_estimate,
                         clone_shape, add_table, place_in_cell,
                         bundled_template, BUNDLED_INDICES)
import mathmark

TEMPLATE = bundled_template()
OUTPUT = "组会汇报-仿射偏压控制.pptx"
MATH = Path("math")
FIG = Path("figures_png")
COVER, TOC = BUNDLED_INDICES["cover"], BUNDLED_INDICES["toc"]
SKELETON, CLOSING = BUNDLED_INDICES["skeleton"], BUNDLED_INDICES["closing"]

CONTENT_L, CONTENT_W = 0.433, 9.292
BANNER_W, BANNER_H, BANNER_TOP = 6.60, 0.96, -0.226
BOTTOM = 7.28
EM = 14 / 72.0
LINE_H = 0.34
CPL = 46
TEXT_BASELINE = 0.208
SYM_DESCENT = 0.042
_MAN = {f["id"]: f for f in json.load(open(MATH / "manifest.json", encoding="utf-8"))["formulas"]}


def nat(fid):
    f = _MAN[fid]
    return f["width_px"] / f["dpi"], f["height_px"] / f["dpi"]


def _load_cjk_font():
    import glob
    from PIL import ImageFont
    for pat in ("C:/Windows/Fonts/msyh.ttc", "C:/Windows/Fonts/msyh.ttf",
                os.path.expanduser("~/Library/Fonts/msyh.ttc"),
                "/usr/share/fonts/**/NotoSansCJK*.*"):
        for path in glob.glob(pat, recursive=True):
            try:
                return ImageFont.truetype(path, 1000)
            except Exception:
                continue
    return None


_FONT = _load_cjk_font()
if _FONT is not None:
    _SPACE_EM = (_FONT.getlength("　") / 1000.0) or 1.0
    def _text_em(ch): return _FONT.getlength(ch) / 1000.0
else:
    _SPACE_EM = 1.0
    def _text_em(ch):
        o = ord(ch)
        if o >= 0x2E80 or ch in "—–→←↔“”‘’…": return 1.0
        if ch == " ": return 0.30
        return 0.62 if "A" <= ch <= "Z" else 0.5


def crop_white(path):
    im = Image.open(path).convert("RGB")
    bbox = ImageChops.difference(im, Image.new("RGB", im.size, "white")).getbbox()
    if bbox: im = im.crop(bbox)
    out = tempfile.mktemp(suffix=".png")
    im.save(out)
    return out


def banner_and_content(slide):
    boxes = [s for s in slide.shapes if s.has_text_frame]
    return min(boxes, key=lambda s: s.top), max(boxes, key=lambda s: s.top)


def H(t): return {"text": t, "size": 20, "bold": True}
def S(t): return {"text": t, "size": 14, "bold": True}
def B(t): return {"text": t, "size": 14, "bold": False}
def R(t): return {"text": t, "size": 14, "bold": False, "color": "C71F2C"}


class Builder:
    def __init__(self, deck):
        self.slide = deck.duplicate(SKELETON)
        self.banner_box, self.tmpl = banner_and_content(self.slide)
        self.y = 1.073
        self._first = True

    def banner(self, text):
        fill_textframe(self.banner_box.text_frame, [{"text": text}])
        self.banner_box.width = Inches(BANNER_W)
        self.banner_box.height = Inches(BANNER_H)
        self.banner_box.top = Inches(BANNER_TOP)
        tf = self.banner_box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    def text(self, items, gap=0.07):
        run = []
        def flush():
            if not run: return
            box = self.tmpl if self._first else clone_shape(self.slide, self.tmpl)
            self._first = False
            box.left, box.top, box.width = Inches(CONTENT_L), Inches(self.y), Inches(CONTENT_W)
            fill_textframe(box.text_frame, list(run))
            self.y = text_bottom_estimate(box) + gap
            run.clear()
        for it in items:
            if mathmark.DELIM.search(it.get("text", "")):
                flush(); self.mb(it["text"], gap=gap)
            else:
                run.append(it)
        flush()

    def mb(self, s, gap=0.07, lead=0.0):
        self.explain(mathmark.parse(s), gap=gap, lead=lead)

    def formula(self, eq_id, gap=0.12, lead=0.08):
        w, h = nat(eq_id)
        self.y += lead
        self.slide.shapes.add_picture(str(MATH / f"{eq_id}.png"), Inches((10.0 - w) / 2),
                                      Inches(self.y), width=Inches(w), height=Inches(h))
        self.y += h + gap

    def explain(self, segments, gap=0.10, lead=0.04):
        self.y += lead
        sp = _SPACE_EM
        atoms = []
        for kind, val in segments:
            if kind == "t":
                for ch in val: atoms.append(("t", ch, _text_em(ch)))
            else:
                w, h = nat(val)
                lx = _MAN[val]["latex"]
                desc = any(t in lx for t in ("_", "varphi", "kappa", "J", "rho",
                                             "beta", "gamma", "ell", "sigma", "eta", "mu"))
                nsp = max(1, math.ceil(w / (sp * EM)))
                atoms.append(("m", val, nsp, w, h, desc, nsp * sp))
        lines, col = [[]], 0.0
        for a in atoms:
            aw = a[2] if a[0] == "t" else a[6]
            if col + aw > CPL and lines[-1]:
                lines.append([]); col = 0.0
            lines[-1].append((a, col)); col += aw
        for li, line in enumerate(lines):
            ly = self.y + li * LINE_H
            string = "".join(a[1] if a[0] == "t" else "　" * a[2] for a, _ in line)
            box = clone_shape(self.slide, self.tmpl)
            box.left, box.top = Inches(CONTENT_L), Inches(ly)
            box.width, box.height = Inches(CONTENT_W), Inches(LINE_H)
            tf = box.text_frame
            tf.word_wrap = False
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            fill_textframe(tf, [{"text": string, "size": 14, "bold": False}])
            tf.paragraphs[0].line_spacing = 1.0
            for a, start in line:
                if a[0] == "m":
                    _, sym, nsp, w, h, desc, awem = a
                    gx = CONTENT_L + start * EM + (awem * EM - w) / 2
                    gy = ly + TEXT_BASELINE + (SYM_DESCENT if desc else 0.0) - h
                    self.slide.shapes.add_picture(str(MATH / f"{sym}.png"), Inches(gx),
                                                  Inches(gy), width=Inches(w), height=Inches(h))
        self.y += len(lines) * LINE_H + gap

    def figure(self, img, width=9.0, gap=0.1, lead=0.06, bottom=None):
        bottom = bottom or BOTTOM
        img = crop_white(str(img))
        self.y += lead
        wpx, hpx = Image.open(img).size
        ar = wpx / hpx
        w = min(width, 9.5); h = w / ar
        if self.y + h > bottom:
            h = bottom - self.y; w = h * ar
        self.slide.shapes.add_picture(img, Inches((10.0 - w) / 2), Inches(self.y),
                                      width=Inches(w), height=Inches(h))
        self.y += h + gap

    def figure_pair(self, img1, img2, height=3.2, gap=0.12, lead=0.06, inner=0.2):
        self.y += lead
        paths = [crop_white(str(img1)), crop_white(str(img2))]
        ars = [(lambda s: s[0] / s[1])(Image.open(p).size) for p in paths]
        if self.y + height > BOTTOM: height = BOTTOM - self.y
        ws = [height * ar for ar in ars]
        total = sum(ws) + inner
        if total > 9.4:
            s = 9.4 / total; height *= s; ws = [w * s for w in ws]; total = sum(ws) + inner
        x = (10.0 - total) / 2
        for p, w in zip(paths, ws):
            self.slide.shapes.add_picture(p, Inches(x), Inches(self.y), width=Inches(w), height=Inches(height))
            x += w + inner
        self.y += height + gap

    def table(self, data, *, col_widths, row_height=0.45, width=None,
              formula_cells=None, gap=0.12, lead=0.08, font_size=12):
        width = width or sum(col_widths)
        self.y += lead
        gf = add_table(self.slide, data, left=(10.0 - width) / 2, top=self.y,
                       width=width, height=row_height * len(data),
                       col_widths=col_widths, row_height=row_height, font_size=font_size)
        for (r, c), fid in (formula_cells or {}).items():
            w, h = nat(fid)
            place_in_cell(self.slide, gf, r, c, MATH / f"{fid}.png", w, h)
        self.y += row_height * len(data) + gap
        return gf


deck = TemplateDeck(TEMPLATE, work_path=OUTPUT)

# ============================ COVER ============================
cover = deck.duplicate(COVER)
for shape in cover.shapes:
    if shape.has_text_frame and "时间" in shape.text_frame.text:
        fill_textframe(shape.text_frame, [
            {"text": "汇报时间：2026.06.17"},
            {"text": "汇报人：　　（请填写）"},
        ])

# ============================ TOC ============================
toc = deck.duplicate(TOC)
fill_textframe(toc.shapes[0].text_frame, [
    H("汇报内容："),
    B("1. 研究背景与问题：任意点偏压控制的结构性困难"),
    B("2. 信号模型与传统导频锁定"),
    B("3. 单 MZM 的精确仿射定理"),
    B("4. 可辨识性、标定与一致增益闭环"),
    B("5. 噪声与导频深度设计"),
    B("6. 推广至 DPMZM：三维环面上的几何"),
    B("7. 可观测性与 IMD 通道的必要性"),
    B("8. 三维环面标定与高斯–牛顿解调"),
    B("9. 偏压控制算法设计"),
    B("10. 数值验证"),
    B("11. 实验验证计划（待实测）"),
    B("12. 讨论、展望与结论"),
])

# ==================================================================
# 1. 研究背景与问题
# ==================================================================
BN1 = "研究背景与问题"

p = Builder(deck); p.banner(BN1)
p.text([
    H("选题：MZM/DPMZM 任意点偏压控制的精确仿射框架"),
    S("核心命题"),
    B("只要接收链与解调过程保持线性，偏置相位与导频时间波形可在"),
    B("传递函数中分离；非理想因素改变的是仿射参数，而非相位特征形式。"),
    S("由此把任意点偏压控制写成"),
    B("“标定模型参数 → 反演相位”的两步问题，对单 MZM 与 DPMZM 统一成立。"),
])
p.figure(FIG / "fig_arch.png", width=6.0)

p = Builder(deck); p.banner(BN1)
p.text([
    H("背景：偏压慢漂移与导频闭环"),
    S("漂移来源"),
    B("LiNbO₃ / InP / 薄膜铌酸锂 MZM 受温度、电荷积累、光折变影响，"),
    B("偏置点在业务运行期间缓慢漂移，需无中断、可在线持续校正的手段。"),
    S("传统导频（dither）闭环"),
    B("小幅低频导频下：一次谐波在零/峰值点过零，二次谐波在正交点过零，"),
    B("由此形成可实现的误差信号 —— 但最容易工作在正交/零/峰值等特殊点。"),
    S("需求"),
    B("微波光子移相/矢量调制、SSB 权重、模拟光计算要求偏置稳在任意设定点。"),
])

p = Builder(deck); p.banner(BN1)
p.text([
    H("任意点控制的三个结构性问题"),
    B(r"传统任意点方案多从特殊点机制外推（«\varphi_b» 信息被压成单一标量谐波）："),
    S("① 死区"),
    B("部分相位附近环路增益趋零，无控制力。"),
    S("② 分支歧义"),
    B(r"«\varphi\leftrightarrow\pi-\varphi» 不可区分（«\sin\varphi» 不变）。"),
    S("③ 增益标定误差"),
    B("标称增益失配直接转化为静态锁定偏差，并被死区因子放大。"),
    R("前两者是拓扑必然，第三者是可辨识性代价；恢复二维观测是同时消除三者的最小代价。"),
])

p = Builder(deck); p.banner(BN1)
p.text([
    H("本文主要贡献"),
    B(r"① 单 MZM 精确仿射定理：«\mathbf{z}=A\mathbf{u}(\varphi_b)+\mathbf{b}+\mathbf{n}»，给出闭式表达。"),
    B("② 统一解释正交/零/峰锁定、幅值匹配、谐波比值法及其静差。"),
    B("③ 自标定任意点控制：椭圆辨识 + 规范固定 + atan2 解调，闭环增益与目标点无关。"),
    B(r"④ DPMZM 推广：状态空间提升为 «\mathbb{T}^3»，12 维三角特征映射保持仿射性。"),
    B("⑤ IMD 通道可观测性解释：标准 QPSK 点秩退化，加互调通道恢复父偏置可观测。"),
    B("⑥ 面向实现的标定与解调算法：每控制周期算力 < 10³ 次乘加。"),
    R("仿真：任意点静差降至毫弧度量级，较传统多环基线降低一至两个数量级。"),
])

# ==================================================================
# 2. 信号模型与传统导频锁定
# ==================================================================
BN2 = "信号模型与传统锁定"

p = Builder(deck); p.banner(BN2)
p.text([H("MZM 模型与解调泛函"), S("推挽单驱 MZM 输出强度")])
p.formula("mzm")
p.text([B(r"式中 «\varphi(t)=\varphi_b+\xi(t)»，导频 «\xi(t)=m\sin\omega t»，«m=\pi V_d/V_\pi»。")])
p.text([S("锁相解调建模为线性泛函（输入波形 → 输出标量）")])
p.formula("lockin")
p.mb(r"式中上横线为低通时间平均，«\theta_1,\theta_2» 为参考相位误差（含群延迟）。")
p.mb(r"观测量 «Y=G_1\mathcal{L}_1[i]»（H1 通道），«X=G_2\mathcal{L}_2[i]»（H2 通道）。")

p = Builder(deck); p.banner(BN2)
p.text([H("谐波结构：Jacobi–Anger 展开"),
        B("偶次谐波携带 cos φb，奇次携带 sin φb，直流携带 1+ηJ₀(m)cos φb：")])
p.table(
    [["频率", "光电流分量", "φb 依赖", "传统用途"],
     ["直流", "", "", "功率监测；规范基准"],
     ["ω", "", "", "零/峰值点锁定"],
     ["2ω", "", "", "正交点锁定"],
     ["3ω", "", "", "（少用）"]],
    col_widths=[1.1, 3.6, 1.6, 2.9], row_height=0.5,
    formula_cells={(1, 1): "hm_dc", (2, 1): "hm_w", (3, 1): "hm_2w", (4, 1): "hm_3w",
                   (1, 2): "dep_c", (2, 2): "dep_s", (3, 2): "dep_c", (4, 2): "dep_s"})
p.text([S("理想极限（θ=0、正弦 ξ、无串扰）")])
p.formula("ideal")

p = Builder(deck); p.banner(BN2)
p.text([H("特殊点锁定及其静差"),
        B("每种特殊点方案只读取 u(φb)=(cos φb, sin φb) 的一个坐标，"),
        B("且只工作在过零点邻域：斜率最大、误差过零（增益标定相消）、另一坐标判符号。"),
        S("正交点 / 零·峰值点静差（仿射记号线性化）")])
p.formula("sQN")
p.mb(r"式中 «A_{ij}»、«b_X»、«b_Y» 为仿射元素；静差即“非对角/偏置元素 ÷ 对角斜率”。")
p.text([R("一旦离开过零点，斜率最大、增益标定相消、符号判别三项便利同时失效。")])

p = Builder(deck); p.banner(BN2)
p.text([H("通往任意点的传统路线"), S("幅值匹配：驱动 Y → 标称参考值")])
p.formula("sAM")
p.mb(r"局部增益 «\propto\eta\kappa_1|\cos\varphi^\ast|» 在正交点附近消失（死区）；"
     r"«Y» 在 «\varphi\leftrightarrow\pi-\varphi» 下不变（分支歧义）。")
p.text([S("谐波比值法：消去公共功率标量")])
p.formula("ratio")
p.mb(r"隐含假设 «A» 对角且 «\mathbf{b}=0»；偏置破坏余切律，«Y\to0» 邻域噪声按 «1/Y^2» 放大。")
p.text([R("一次除法只能消去一个公共标量；真实失配是六参数的仿射群作用。")])

p = Builder(deck); p.banner(BN2)
p.text([H("拓扑障碍（命题 1）"),
        S("单标量谐波误差函数的内在缺陷"),
        B("由单个标量谐波特征 g(φb) 构造的任何误差函数都是光滑 2π 周期函数："),
        B("· Rolle 定理 ⇒ 每周期内必有临界点 g′(φ†)=0，该处环路增益消失；"),
        B("· 周期函数在一个周期上不可能单射 ⇒ 分支歧义不可避免。"),
        S("二维特征的对照"),
        B(r"映射 «\varphi_b\mapsto\mathbf{u}(\varphi_b)» 是浸入 «\mathbb{S}^1\to\mathbb{R}^2»，"),
        B(r"满足 «\|\mathbf{u}'\|\equiv1»（处处有增益）且在 «[0,2\pi)» 上单射（无歧义）。"),
        R("二维化是消除死区与歧义的最小代价；但物理给出的是 u 的未知形变像。"),
])

# ==================================================================
# 3. 单 MZM 精确仿射定理
# ==================================================================
BN3 = "单 MZM 精确仿射定理"

p = Builder(deck); p.banner(BN3)
p.text([H("相位因子分解（引理 1）"),
        B("对任意波形 ξ(t)，下式精确成立——偏置相位与时间依赖完全因子化：")])
p.formula("lemma")
p.text([S("含义"),
        B("导频波形决定线性映射（如何把单位圆拉成椭圆），"),
        B("偏置相位决定单位圆上的位置。两者在传递函数中天然分离。"),
        R("这是整个仿射框架的代数起点：不必小信号近似，对任意周期导频成立。")])

p = Builder(deck); p.banner(BN3)
p.text([H("精确仿射定理（定理 1）"),
        B("任意周期导频 + 任意线性接收链 + 线性解调泛函下，下式精确成立：")])
p.formula("affine")
p.text([S("仿射矩阵 A")])
p.formula("Amat")
p.text([S("偏置 b 与噪声 n")])
p.formula("bnvec")
p.mb(r"常数项 «\rho P_0/2» 被两个泛函湮灭；«O_X,O_Y» 为混频器直流失调。")

p = Builder(deck); p.banner(BN3)
p.text([H("线性非理想 → 仿射参数（不改模型形式）"),
        B("线性非理想因素通过 A、b 改变观测椭圆的位置/尺度/方向；相位仍只经 u(φb) 进入。"),
        S("理想对角元（传统谐波斜率）")])
p.formula("Adiag")
p.mb(r"以含二次谐波失真的导频 «\xi=m\sin\omega t+\varepsilon m\sin(2\omega t+\psi)» 为例：")
p.text([B("至一阶，导频二次谐波失真主要生成 A 的非对角耦合（详见附录）；")])
p.mb(r"理想极限下 «A\to\operatorname{diag}(\eta\kappa_2,-\eta\kappa_1)»。")
p.mb(r"分界的工程意义：与导频相干的确定性干扰进 «\mathbf{b}»（可标定补偿），非相干噪声进 «\mathbf{n}»（只能统计抑制）。")

p = Builder(deck); p.banner(BN3)
p.text([H("观测椭圆与模型边界"),
        S("几何：无噪观测落在椭圆上")])
p.formula("ellipse")
p.mb(r"«M\triangleq(AA^{\mathsf{T}})^{-1}\succ0» 为定义椭圆的正定二次型；中心为 «\mathbf{b}»。")
p.text([S("模型边界（需单独处理的非仿射效应）"),
        B("· 光电探测/跨阻非线性 ⇒ cos 2φb 型非仿射残差（小信号下二阶）；"),
        B("· 部分 InP 器件偏压相关吸收（电吸收/残余幅度调制）⇒ 乘性因子 T(φb)；"),
        B("· 导频深度 m 的慢漂移 ⇒ A 缓变，需重定标处理。")])

# ==================================================================
# 4. 可辨识性、标定与一致增益闭环
# ==================================================================
BN4 = "标定与一致增益闭环"

p = Builder(deck); p.banner(BN4)
p.text([H("规范自由度：S¹ 的稳定子 O(2)"),
        B("椭圆有 5 个参数，而 (A, b) 有 6 个自由度，差额来自单位圆的稳定子："),
        S("稳定子"),
        B(r"对任意 «R\in O(2)»，«A\mapsto AR» 不改变椭圆像集"),
        B("（u 跑遍单位圆时 Ru 也跑遍单位圆）。"),
        S("后果"),
        B(r"单靠椭圆几何只能把 A 辨识到右乘正交阵的等价类，相位仍有 «\hat\varphi_b\mapsto\pm\hat\varphi_b+\varphi_c» 歧义。"),
        R("规范固定需额外信息：直流通道提供相位原点，扫描方向提供定向。")])

p = Builder(deck); p.banner(BN4)
p.text([H("四步标定算法"),
        B("(i) 对偏压作满周期慢扫，记录 {zᵢ, īᵢ}；"),
        B("(ii) 中心化/归一后代数最小二乘拟合二次曲线 ⇒ 中心 b̂、M≻0（高噪用 Fitzgibbon）；"),
        B("(iii) 由驻点解中心、取对称平方根 M^½，由轨道环绕方向固定反射；"),
        B("(iv) 旋转规范由整条直流序列对 (1, cos θ, sin θ) 回归固定 φ_c。"),
        S("解调器")])
p.formula("demod")
p.mb(r"直流回归把标定引入的解调偏差中位从 «{\sim}42» mrad 降至 «{\sim}2» mrad（«\sigma=0.005»）。")

p = Builder(deck); p.banner(BN4)
p.text([H("标定示例"),
        B("m=1.2、σ=0.005、N=360，含增益失配、10° 非正交、交叉耦合 ε=0.10 与偏置："),
        B("(a) 原始观测平面：拟合椭圆、中心 b̂、φ=π 基准点；(b) 回拉样本 û=B(z−b̂) 落于单位圆。")])
p.figure(FIG / "fig_ellipse.png", width=8.2)

p = Builder(deck); p.banner(BN4)
p.text([H("一致增益闭环与方案对照"),
        B("解调后相位误差与真实偏差直接对应，不依赖目标点；配合常数执行斜率 π/Vπ，"),
        B("同一组 PI(D) 参数可用于整个偏置周期。")])
p.table(
    [["方案", "可锁点集", "增益", "静差", "歧义性"],
     ["H2 过零（正交）", "{π/2,3π/2}", "ηκ₂", "", "半周期"],
     ["H1 过零（零/峰）", "{0,π}", "ηκ₁", "", "半周期"],
     ["幅值匹配", "死区之外", "ηκ₁|cosφ*|", "式(AM)", "φ↔π−φ"],
     ["谐波比值", "Y 非小", "随目标点", "b 主导", "双象限"],
     ["仿射法（本文）", "[0,2π)", "一致", "0（模型内）", "无"]],
    col_widths=[2.1, 1.8, 1.7, 1.9, 1.5], row_height=0.46,
    formula_cells={(1, 3): "sQ", (2, 3): "sN"})
p.text([R("仿射法零静差指模型内不含死区、分支歧义与增益标定偏差，与漂移跟踪误差相互独立。")])

p = Builder(deck); p.banner(BN4)
p.text([H("单 MZM 自标定 + 闭环算法"),
        S("标定阶段（上电执行一次）"),
        B("① 慢扫偏压覆盖 [0,2π)，记录 {z, ī}；"),
        B("② 代数最小二乘拟合锥曲线 ⇒ 中心 b̂、M=(AAᵀ)⁻¹；"),
        B("③ B₀=M^½，由环绕方向定反射 F；④ 直流回归得相位原点 φ_c，B=R(−φ_c)F B₀。"),
        S("运行阶段（每控制周期）"),
        B("① 读锁相 z=(Y,X)；② û=B(z−b̂)，φ̂=atan2(û_y,û_x)；"),
        B("③ e=wrap(φ̂−φ*)；④ PI 更新 Vb ← Vb − G·e·Vπ/π 写 DAC；"),
        B("⑤ 残差 ρ=|‖û‖−1| 持续越限 ⇒ 触发重定标。"),
        R("atan2 给出与工作点无关的相位误差，是相对传统方法的核心优势。")])

# ==================================================================
# 5. 噪声与导频深度设计
# ==================================================================
BN5 = "噪声与导频深度设计"

p = Builder(deck); p.banner(BN5)
p.text([H("相位噪声与条件数"), S("沿切向一阶摄动")])
p.formula("var")
p.mb(r"几何含义：相位只沿椭圆切向移动，相位噪声 = 测量噪声 «\sigma» ÷ 椭圆在该方向的拉伸量。")
p.text([S("条件数即椭圆扁度")])
p.formula("kappa")
p.mb(r"理想链路 «A=\operatorname{diag}(\eta\kappa_2,-\eta\kappa_1)»，故 «\kappa(A)=J_1(m)/J_2(m)\to4/m»。")
p.mb(r"最坏情形由最小奇异值 «\sigma_{\min}(A)» 封顶；插入式估计当 «\kappa(A)» 适中时与 CRB 仅差二阶。")

p = Builder(deck); p.banner(BN5)
p.text([H("导频深度选择：两难权衡"),
        B("· m 过小：J₂≈m²/8 很小，H2 通道幅度被噪声淹没，κ(A) 按 1/J₂≈8/m² 急剧放大；"),
        B("· m 过大：导频本身对业务信号的调制损伤随 m² 增大，并激起更高次谐波。"),
        B("折中区即设计取值，本文取 m=1.2（理想比值 3.13，含失配真值矩阵 2.59）。")])
p.figure(FIG / "fig_bessel.png", width=6.6)
p.text([R("κ(A) 可作为选择导频深度的单一量化指标。")])

# ==================================================================
# 6. 推广至 DPMZM：T³ 几何
# ==================================================================
BN6 = "推广至 DPMZM：T³ 几何"

p = Builder(deck); p.banner(BN6)
p.text([H("场模型与半角结构"),
        B("DPMZM 由两子 MZM 嵌套于父干涉仪，偏置相位 φ=(φ₁,φ₂,φ₃)。无 RF 数据时：")])
p.formula("dpfield")
p.formula("dppower")
p.mb(r"前两项为两子 MZM 自项；末项来自父干涉仪交叉干涉，含子相位的半角 «\cos(\varphi_i/2)»。")
p.text([R("半角结构是 DPMZM 相对单 MZM 的关键差异，带来三个直接后果。")])

p = Builder(deck); p.banner(BN6)
p.text([H("半角结构的三个直接后果"),
        S("① 4π 周期"),
        B("cos(φ₁/2) 周期为 4π：φ₁→φ₁+2π 时子场与交叉项同时变号，输出可区分；"),
        B("子轴标定扫描应覆盖 4π（电压 4Vπ），是单 MZM 常规范围的两倍。"),
        S("② Klein 等价")])
p.formula("klein")
p.mb(r"二者生成 «V_4=\{\mathrm{id},T_1,T_2,T_1T_2\}»；构型空间为商 «\widetilde{\mathbb{T}}^3/V_4»。")
p.text([S("③ 任意点运行"),
        B("标准 QPSK 偏置为 (π,π,±π/2)（双子零点、父正交、载波抑制）；"),
        B("任意点运行（矢量调制、SSB 权重）把三轴停在一般位置。")])

p = Builder(deck); p.banner(BN6)
p.text([H("强度图与 Klein 商"),
        B("强度 P(φ₁,φ₂; φ₃=π/2) 在 [0,4π)² 区片：标准 QPSK 点（实心）与其三个 Klein"),
        B("像（空心）具有相同强度，对应 V₄ 商结构；对 QPSK 仅是星座重标记与全局相位。")])
p.figure(FIG / "fig_torus.png", width=5.2)

p = Builder(deck); p.banner(BN6)
p.text([H("12 维三角特征映射与仿射定理"),
        B("对三相位各用一次因子分解（半角版逐字成立）并三线性展开交叉项，"),
        B("光电流的 φ 依赖只能由十二个三角单项式构成：")])
p.formula("feature")
p.text([S("DPMZM 精确仿射性（定理 2）")])
p.formula("dpaffine")
p.mb(r"理想链路下 «A=A_0» 由九通道系数给出，且高度稀疏（每行一至二个非零元）。")
p.text([B("做法：先固定特征映射 Φ，再用线性回归辨识 A、b 以吸收线性链路失配。")])

p = Builder(deck); p.banner(BN6)
p.text([H("结构对应：单 MZM ↔ DPMZM")])
p.table(
    [["要素", "单 MZM", "DPMZM"],
     ["状态流形", "S¹", "T³（构型 T̃³/V₄）"],
     ["特征映射", "u∈R²（恒等）", "Φ∈R¹²（含半角）"],
     ["轴向周期", "2π", "子轴 4π，父轴 2π"],
     ["通道数", "2（ω, 2ω）", "≥7，IMD 必需"],
     ["曲面拟合", "椭圆（二次簇）", "准周期扫描回归"],
     ["规范群", "O(2)", "T³⋊(反射×V₄)"],
     ["解调", "atan2（闭式）", "热启动高斯–牛顿"],
     ["不可观修复", "H1/H2 联用", "IMD 通道"]],
    col_widths=[2.2, 3.0, 3.8], row_height=0.46)

# ==================================================================
# 7. 可观测性与 IMD 通道
# ==================================================================
BN7 = "可观测性与 IMD 通道"

p = Builder(deck); p.banner(BN7)
p.text([H("观测雅可比与秩条件"),
        B("解调局部可行性由观测雅可比的秩条件刻画（命题 1 浸入论证的三维推广）：")])
p.formula("jac")
p.text([S("命题 2（标准点不可观及其奇异轨迹）"),
        B("取六个谐波通道 Y₁₂₃, X₁₂₃ 与理想 A₀，使父轴不可观的集合为 {c₁=c₂=0}，"),
        B("即标准 QPSK 工作点的 Klein 轨道，且对任意 φ₃ 成立。"),
        B("加入 IMD 通道 Z₋ 后，单项式 s₁s₂C₃ 的 φ₃ 偏导在标准点为 ∓1，恢复满秩。"),
        R("六谐波通道在标准 QPSK 点缺少父偏置方向信息；IMD 音检测父偏置 = 恢复雅可比秩。")])

p = Builder(deck); p.banner(BN7)
p.text([H("可观测性热图"),
        B("log₁₀σmin(𝒥) 于 (φ₁,φ₂)∈[0,4π)²：暗紫=近奇异（盲区），亮黄=可观测。")])
p.figure(FIG / "fig_obs.png", width=9.4, bottom=6.7)
p.text([B("(a) 六谐波/φ₃=π/2：低秩区落在 Klein 轨道（4 个暗斑）；(b) 加 IMD 恢复；"),
        B("(c) 九通道/φ₃=0：残余奇点 (π,π,0)（输出全暗，σmin≈1.1×10⁻¹⁷）。")])

# ==================================================================
# 8. T³ 标定与高斯–牛顿
# ==================================================================
BN8 = "T³ 标定与高斯–牛顿"

p = Builder(deck); p.banner(BN8)
p.text([H("准周期扫描线性回归"),
        B("Φ(T³) 非二次簇，无闭式椭圆拟合；但观测对 Φ 仍线性，可用普通最小二乘："),])
p.formula("regress")
p.mb(r"一次 «13\times13» Gram 分解服务所有 K 个通道；取增量比无理的准周期轨线 «\varphi_i(k)=k\Delta_i»，单条轨线即可覆盖 «\mathbb{T}^3»。")
p.figure(FIG / "fig_ahat.png", width=5.6)
p.text([B("回归所得 Â（色块）对照理想 A₀ 稀疏指纹（方框）；框外色块为注入失配 δA=0.12。")])

p = Builder(deck); p.banner(BN8)
p.text([H("热启动高斯–牛顿解调"),
        B("圆嵌入有全局闭式逆，Φ(T³) 没有；但偏压控制是跟踪问题，"),
        B("可用上一周期估计热启动，执行阻尼高斯–牛顿步：")])
p.formula("gn")
p.mb(r"本文仿真设置下两次迭代即达测量噪声地板；收敛域由 «\sigma_{\min}(\mathcal{J})» 与嵌入曲率共同决定。")
p.text([B("借助 Â 稀疏性，每控制周期约数百次乘加 + 一次 3×3 求解，适合 Cortex-M 级控制器。")])

p = Builder(deck); p.banner(BN8)
p.text([H("父轴噪声与对照基线"),
        S("估计协方差为 Cramér–Rao 型")])
p.formula("covcrb")
p.text([S("双子零点附近父相位方差")])
p.formula("parentnoise")
p.mb(r"高于子轴地板 «(8/m)^2\sigma^2»，比值 «\propto1/m^2»；浅导频（«m\approx0.13»）时可达一个量级以上。")
p.text([S("对照基线"),
        B("文献式三独立环：误差量取 (Y₁,Y₂,Z₋)，参考由理想模型算，斜率取理想雅可比对角元；"),
        B("保留理想耦合项，但不做矩阵反演、不辨识 A≠A₀ 与 b≠0 的线性失配。")])

# ==================================================================
# 9. 偏压控制算法设计
# ==================================================================
BN9 = "偏压控制算法设计"

p = Builder(deck); p.banner(BN9)
p.text([H("两阶段监督状态机"),
        B("标定阶段：扫描 → 拟合与规范固定 → 自检；运行阶段每周期：锁相读出 → 解调 →"),
        B("PI → 残差监测；残差越限持续 M 周期即触发微幅弧重定标并返回标定通道。")])
p.figure(FIG / "fig_flow.png", width=9.0, bottom=6.9)
p.text([B("单 MZM 与 DPMZM 共用骨架，区别仅在拟合方式与解调算子（atan2 或 高斯–牛顿）。")])

p = Builder(deck); p.banner(BN9)
p.text([H("整定与实现"),
        S("阈值整定"),
        B("ρth 由标定末残差分布给出，配 EWMA（λ≈0.05）与连续 M（≈20–30）周期判决；"),
        B("检测延迟量级约 M+1/λ 个控制周期。"),
        S("重定标弧"),
        B("±0.3 rad 量级微幅弧重新激励局部可辨识子空间（DPMZM 优先激励 s₁s₂ 子空间）。"),
        S("定点化"),
        B("atan2 用 CORDIC/查表；3×3 求解用 Cholesky 展开为定长乘加序列。"),
        S("参数存储量"),
        B("单 MZM 六个标量；DPMZM 为 9×12+9 = 117 个标量。全周期代价 < 10³ 次乘加。")])

# ==================================================================
# 10. 数值验证
# ==================================================================
BN10 = "数值验证"

p = Builder(deck); p.banner(BN10)
p.text([H("仿真设置与参数"),
        B("受扰真值 A真值=A₀+δA·ΔA、b真值=δb·Δb，δA=0.12（相对 Frobenius 尺度）；"),
        B("叠加逐通道高斯噪声与共同漂移序列（随机游走+慢正弦），对比控制器施同源漂移。")])
p.table(
    [["参数组", "取值"],
     ["单 MZM", "m=1.2、gX=1.18、gY=0.88、δ=10°、ε=0.10"],
     ["", "b=(0.030, −0.022)、σ=0.005、N=360"],
     ["DPMZM", "m₁=m₂=m₃=1.2、σ=0.002、N=3000"],
     ["导频频率", "ω₁:ω₂:ω₃ = 1 : 1.37 : 1.93（互调不碰撞）"]],
    col_widths=[2.0, 7.0], row_height=0.5)
p.mb(r"日常电压量经 «m=\pi V_d/V_\pi»、«\varphi_b=\pi V_b/V_\pi+\varphi_0» 换成相位量，或并入归一化 «\kappa_n»。")

p = Builder(deck); p.banner(BN10)
p.text([H("主要结果（单 MZM）")])
p.table(
    [["量", "结果"],
     ["零噪声解调误差（全周期）", "0.0000 mrad（数值精度）"],
     ["零噪声 Â 对真值 A", "Frobenius 误差 0.00%"],
     ["含噪标定解调偏差（σ=0.005）", "中位 2.35、P95 6.25 mrad"],
     ["规范固定：直流回归 对 argmin", "~2 对 ~42 mrad"],
     ["闭环（φ*=1.9）", "仿射 12.7 / H1 匹配 1315 mrad（rms）"],
     ["m=1.2 处 κ(A)（受扰真值）", "2.59（理想比值 3.13）"]],
    col_widths=[4.6, 4.6], row_height=0.55)
p.text([R("式 (affine) 在仿真模型内为精确结构，而非小信号近似。")])

p = Builder(deck); p.banner(BN10)
p.text([H("主要结果（DPMZM）")])
p.table(
    [["量", "结果"],
     ["DΦ 对有限差分", "最大 7.7×10⁻¹⁰"],
     ["Klein 不变性 ‖Φ∘T₁−Φ‖", "2.5×10⁻¹⁵"],
     ["零噪声辨识 / GN 恢复", "2.4×10⁻¹⁴% / 5×10⁻¹³ mrad"],
     ["含噪辨识（σ=0.002）/ 解调偏差", "0.17% / 中位 2.1、P95 4.3 mrad"],
     ["σmin：六 / 九通道（标准点）", "≈0（机器精度）/ 5.5×10⁻²"],
     ["σmin：九通道于 (π,π,0)", "1.1×10⁻¹⁷"],
     ["闭环（任意目标点）", "GN 17.2 / 三环 358 mrad（rms）"],
     ["闭环（标准目标点）", "GN 18.5 / 三环 397 mrad（rms）"]],
    col_widths=[4.6, 4.6], row_height=0.46)

p = Builder(deck); p.banner(BN10)
p.text([H("单 MZM 闭环对比"),
        B("任意目标点 φ*=1.9 rad，同源漂移，σ=0.005（对数误差轴）：")])
p.figure(FIG / "fig_mzmloop.png", width=6.6)
p.text([B("仿射解调保持 12.7 mrad rms；H1 幅值匹配锁入错误分支，静差 ~1.3 rad。")])

p = Builder(deck); p.banner(BN10)
p.text([H("任意点性能随目标相位"),
        B("(a) 传统两类方案闭式静差在各自死区附近放大，仿射环路实测 rms 随目标点变化很小；"),
        B("(b) 解调噪声地板：式 (var) 理论曲线、蒙特卡洛实测与上界 σ/σmin(A) 一致。")])
p.figure(FIG / "fig_sweep.png", width=9.4, bottom=6.9)

p = Builder(deck); p.banner(BN10)
p.text([H("DPMZM 闭环跟踪"),
        B("同源漂移，σ=0.002、δA=0.12（对数误差轴）：")])
p.figure(FIG / "fig_dploop.png", width=9.2, bottom=6.6)
p.text([B("(a) 任意目标点：三环基线 ~0.36 rad 静差；(b) 标准 QPSK 点基线仍经二阶小父轴斜率"),
        B("继承可比静差；高斯–牛顿仿射控制器在两处均保持 ~18 mrad rms。")])

p = Builder(deck); p.banner(BN10)
p.text([H("蒙特卡洛静差 与 规范固定对照"),
        B("左：30 个失配实现的 DPMZM 静差分布（三环基线携数百 mrad，任意目标处出现发散）；"),
        B("右：规范固定方案对照——直流回归按 1/√N 收缩，argmin 基准受平坦极小值限制。")])
p.figure_pair(FIG / "fig_mcdp.png", FIG / "fig_gauge.png", height=3.4)

p = Builder(deck); p.banner(BN10)
p.text([H("算法级动态：捕获与阶跃"),
        B("(左) 捕获瞬态：不同目标点误差曲线重合并贴合理论包络 (1−G)ⁿ；"),
        B("(右) 设定点阶跃：各步整定 20–22 周期（至 20 mrad），与目标点无关。")])
p.figure_pair(FIG / "fig_acq.png", FIG / "fig_step.png", height=2.7)
p.mb(r"时间常数 «1/G»：MZM «\approx5.6»、DPMZM «\approx6.2» 个周期。")

p = Builder(deck); p.banner(BN10)
p.text([H("算法级动态：残差触发重定标"),
        B("第 1200 周期注入漂移突变（gX:1.18→0.85 且 bx 跳变）："),
        B("圆残差 EWMA 越过阈值 ρth=0.06，经 31 周期判决延迟触发重定标。")])
p.figure(FIG / "fig_recal.png", width=8.6, bottom=6.7)
p.text([R("失配期锁定误差升至 275 mrad，重定标后恢复至 13.7 mrad（突变前 12.6）。")])

# ==================================================================
# 11. 实验验证计划
# ==================================================================
BN11 = "实验验证计划（待实测）"

p = Builder(deck); p.banner(BN11)
p.text([H("测量链路（单 MZM）"),
        B("TLS1000 激光（1550 nm、+10 dBm）入射 MZM，1:9 耦合器 10% 支路送 PP-10G"),
        B("探测器，电输出进 SDS824X HD 示波器；上位机由 FFT 读观测向量 z，完成仿射反演"),
        B("与 PI，再经 DG922 Pro 产生偏置与导频写回。RF 关断。")])
p.figure(FIG / "fig_exp_mzm.png", width=8.8, bottom=6.7)
p.mb(r"导频 400 mVpp/音，«m=\pi A_p/V_\pi\approx0.13»，«\kappa(A)\approx30»，须以 FFT 平均换信噪比。")

p = Builder(deck); p.banner(BN11)
p.text([H("测量链路（DPMZM）"),
        B("与单 MZM 同构，区别在三路偏置与三音导频（ω₁,ω₂,ω₃），由两台 DG922 Pro"),
        B("（共六通道）产生；示波器 FFT 同时读取六谐波与三个互调通道。")])
p.figure(FIG / "fig_exp_dpmzm.png", width=8.8, bottom=6.6)
p.text([B("规范固定所需直流由另一路慢速探测器获取（PP-10G 交流耦合，低频截止 ~30 kHz）。"),
        R("严禁编造实测数据：本节为计划，占位图与“待测”单元待平台测试后填入。")])

p = Builder(deck); p.banner(BN11)
p.text([H("分阶段流程与判据"),
        B("(0) 平台与 Vπ 标定；(1) 单 MZM 标定 ↔ 图(椭圆)；(2) 任意点锁定 ↔ 图(扫描)；"),
        B("(3) 噪声与导频深度，扫 m 核对 κ(A)；(4) 漂移与重定标，24h 连续运行；"),
        B("(5) DPMZM 子轴 4π 与 Klein；(6) IMD 可观测性恢复；(7) 任意点三轴锁定。"),
        S("判据"),
        B("各实测指标与仿真同量级；任意点 rms 不随目标相位系统性放大；"),
        B("IMD 通道确能在标准点附近恢复父偏置可观测性。")])

p = Builder(deck); p.banner(BN11)
p.text([H("实验测试矩阵（节选）")])
p.table(
    [["指标", "仿真值", "实验值"],
     ["标定自检残差（MZM）", "中位 2.35 mrad", "待测"],
     ["任意点锁定 rms", "12.7 mrad@1.9 rad", "待测"],
     ["噪声地板 κ(A)", "2.59（m=1.2）", "待测（m≈0.13：≈30）"],
     ["DPMZM 辨识 relF", "0.17%", "待测"],
     ["IMD σmin（六→九，标准点）", "≈0 → 5.5×10⁻²", "待测"],
     ["任意点三轴 rms（GN/三环）", "17.2 / 358 mrad", "待测"],
     ["标准点三轴 rms（GN/三环）", "18.5 / 397 mrad", "待测"]],
    col_widths=[3.8, 2.9, 2.5], row_height=0.5)

# ==================================================================
# 12. 讨论、展望与结论
# ==================================================================
BN12 = "讨论、展望与结论"

p = Builder(deck); p.banner(BN12)
p.text([H("适用边界与开放问题"),
        S("持续激励"),
        B("闭环稳态下样本不再充分激励参数空间（辨识型控制共性）；用残差监测+重定标缓解。"),
        S("漂移时间尺度"),
        B("热致漂移秒~分钟、电荷/光折变分钟~小时；取控制周期 1–10 ms，导频带宽远高于漂移谱。"),
        S("DP-QPSK"),
        B("独立监测 PD 时模型块对角，两偏振分别套用本框架；共享 PD 时需扩展特征空间。"),
        S("适用边界"),
        B("接收链非线性、偏压相关吸收、直流基准 O(ε) 规范偏移限制绝对精度。")])

p = Builder(deck); p.banner(BN12)
p.text([H("展望"),
        B("· 器件维度推广：凡传递函数能按偏置相位与导频波形分离者（IQ 嵌套、级联调制器、"),
        B("  薄膜铌酸锂多电极）皆可写出相应 Φ 与仿射模型；"),
        B("· 半参数辨识：把接收链非线性与偏压相关吸收作为已知结构纳入“仿射+非线性”模型；"),
        B("· 在线自适应：子空间受限递推最小二乘 + 受激励调度，把重定标推进为持续自适应；"),
        B("· 导频规划在线优化：以噪声协方差（式 var、Fisher/CRB）为目标，κ(A) 与父轴 IMD 信噪比为指标；"),
        B("· 通道设计准则：由器件结构直接推出恢复满秩所需最小导频/谐波/互调通道集。")])

p = Builder(deck); p.banner(BN12)
p.text([H("结论"),
        B("· 从偏置相位与导频时间波形的因子分解出发，证明 MZM 谐波解调观测是单位圆的"),
        B("  精确仿射像，并推广到 DPMZM 的 12 维三角特征映射；"),
        B("· 特殊点锁定及其静差可解释为仿射观测的低维读出；任意点控制经标定、规范固定、"),
        B("  相位反演实现，闭环增益不再依赖目标点位置；"),
        B("· DPMZM 半角结构给出 4π 子轴周期、Klein 四元群等价与 IMD 通道秩条件解释；"),
        R("· 仿真：毫弧度量级残差，较传统多环基线降低任意点静差一至两个数量级。"),
        B("主要边界：接收链非线性、偏压相关吸收与持续激励，待残差监测/重定标/实验进一步量化。")])

# ============================ CLOSING ============================
deck.duplicate(CLOSING)

# ============================ BACKUP ============================
BNB = "附录与备份"

p = Builder(deck); p.banner(BNB)
p.text([H("附录：导频失真下的首阶元素"),
        B("设 ξ=m sinωt + εm sin(2ωt+ψ)，展开并用投影恒等式得：")])
p.formula("appL")
p.text([S("代入即得 A 的非对角耦合（对角元不变）")])
p.formula("Aoff")
p.text([S("直流曲线整体相移（限制绝对相位精度）")])
p.formula("dphic")

p = Builder(deck); p.banner(BNB)
p.text([H("备份：DPMZM 九通道首阶系数"),
        B("自项展开用全深度 Jₙ(mᵢ)，交叉项用半深度 jₙ(mᵢ/2)；下表为首阶系数×单项式：")])
p.table(
    [["通道", "频率", "首阶系数 × 单项式", "角色"],
     ["Y₁", "ω₁", "", "子 I 锁定"],
     ["X₁", "2ω₁", "", "判别"],
     ["Y₂", "ω₂", "", "子 Q 锁定"],
     ["X₂", "2ω₂", "", "判别"],
     ["Y₃", "ω₃", "", "父（注意 c₁c₂）"],
     ["X₃", "2ω₃", "", "父判别"],
     ["Z₋", "ω₁−ω₂", "", "Kawakami 型"],
     ["Z₁₃", "ω₃−ω₁", "", "本文引入"],
     ["Z₂₃", "ω₃−ω₂", "", "本文引入"]],
    col_widths=[0.9, 1.3, 5.3, 1.7], row_height=0.46, font_size=11,
    formula_cells={(1, 2): "ch_y1", (2, 2): "ch_x1", (3, 2): "ch_y2", (4, 2): "ch_x2",
                   (5, 2): "ch_y3", (6, 2): "ch_x3", (7, 2): "ch_zm",
                   (8, 2): "ch_z13", (9, 2): "ch_z23"})

# ============================ FINISH ============================
deck.keep_only_new()
deck.save()
print(f"saved {OUTPUT}  ({len(deck.prs.slides)} slides)")
